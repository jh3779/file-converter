"""LibreOffice 번들 스모크 — 실제 앱 코드 경로로 DOCX/PPTX→PDF, HWP→PDF, DOCX→HWP, 영상→MP4 전체 파이프라인 검증.

CI(Windows)에서 사용. 기본은 FILECONV_SOFFICE / FILECONV_JAVA / FILECONV_HWP_CLASSPATH
환경변수로 엔진 위치를 지정(사전 빌드 경로 검증용 — dist가 아직 없는 시점).

--frozen-exe <exe경로>를 주면 환경변수를 전혀 쓰지 않고, sys.frozen/sys.executable을
그 경로로 흉내내어 app.bundle.engine_dir()의 실제 프로덕션 자동 탐색 로직
(exe 옆 engine/ 폴더)이 진짜로 동작하는지 검증한다 — 이것이 실사용자가
FileConverter.exe를 실행했을 때 실제로 타는 경로다.

PDF 결과는 파일 크기뿐 아니라 %PDF- 매직바이트 + pdfminer 텍스트 추출로
한글 내용까지 확인한다(빈 PDF·깨진 PDF를 파일 크기만으로 오검출하지 않도록).

실패 시 non-zero로 종료 — 이 스크립트를 호출하는 쪽에서 exit code를 반드시 확인할 것.

사용:
  python scripts/smoke_pdf_pipeline.py <hwp_sample_path> [--skip-hwp] [--frozen-exe <exe경로>]
"""
import shutil
import sys
import tempfile
from pathlib import Path

# Windows CI 콘솔 기본 인코딩(cp1252)은 화살표(→)·한글을 출력하지 못해
# print()가 UnicodeEncodeError로 죽는다 — 실제 변환 결과와 무관한 순수 출력 버그이므로
# 가장 먼저 UTF-8로 강제 재구성한다.
sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")

REPO = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO))

from app import converters  # noqa: E402
from app.converters.base import ConversionError  # noqa: E402
from app.converters import office as office_mod  # noqa: E402
from app.converters import hwp as hwp_mod  # noqa: E402
from app.converters import video as video_mod  # noqa: E402


def check(label: str, cond: bool, detail: str = ""):
    status = "OK" if cond else "FAIL"
    print(f"[{status}] {label}" + (f" — {detail}" if detail else ""))
    if not cond:
        sys.exit(1)


def _font_actually_embedded(path: Path, font_substr: str) -> bool:
    """FontDescriptor 객체를 실제로 파싱해 /FontFile(2/3) 스트림 참조가 있는지 확인한다.

    단순 'b"NotoSansKR" in raw' 문자열 검사는 폰트가 실제 내장되지 않고
    이름만 참조된 경우(/BaseFont 등)도 통과시킬 수 있어 불충분하다 —
    pdfminer로 PDF 객체를 순회해 /Type /FontDescriptor 이면서 /FontName에
    font_substr이 포함된 항목에 /FontFile·/FontFile2·/FontFile3 중 하나가
    실제로 있는지 확인해야 '내장'의 증거가 된다.
    """
    from pdfminer.pdfparser import PDFParser
    from pdfminer.pdfdocument import PDFDocument
    from pdfminer.pdftypes import resolve1, PDFObjRef

    with path.open("rb") as f:
        doc = PDFDocument(PDFParser(f))
        xref = doc.xrefs[0] if doc.xrefs else None
        obj_ids = set()
        for x in doc.xrefs:
            obj_ids.update(x.get_objids())
        for objid in obj_ids:
            try:
                obj = resolve1(doc.getobj(objid))
            except Exception:
                continue
            if not isinstance(obj, dict) or obj.get("Type") is None:
                continue
            type_name = str(obj.get("Type"))
            if "FontDescriptor" not in type_name:
                continue
            font_name = str(obj.get("FontName", ""))
            if font_substr not in font_name:
                continue
            if any(k in obj for k in ("FontFile", "FontFile2", "FontFile3")):
                return True
    return False


def verify_pdf_content(path: Path, must_contain: str, must_embed_font: bool = False) -> None:
    raw = path.read_bytes()
    check(f"{path.name}: %PDF- 매직바이트", raw[:5] == b"%PDF-")
    from pdfminer.high_level import extract_text
    text = extract_text(str(path))
    check(f"{path.name}: 텍스트 내용('{must_contain}') 포함",
          must_contain in text, repr(text[:120]))
    if must_embed_font:
        # 텍스트 추출은 ToUnicode CMap만 보므로 코드값이 맞아도 실제 렌더링
        # 글리프가 깨질 수 있다(DEC-015 재현 당시 직접 확인) — 폰트 이름이
        # 문자열로 등장하는 것만으론 부족하고, FontDescriptor에 실제
        # /FontFile* 스트림이 물려 있는지까지 객체 단위로 확인해야 한다.
        check(f"{path.name}: NotoSansKR 폰트 실제 내장(FontDescriptor/FontFile* 확인)",
              _font_actually_embedded(path, "NotoSansKR"))


def simulate_frozen(exe_path: str):
    """실제 사용자가 FileConverter.exe를 실행했을 때와 동일하게
    sys.frozen/sys.executable을 흉내내, engine_dir() 자동 탐색을 진짜로 태운다."""
    resolved = Path(exe_path).resolve()
    sys.frozen = True
    sys.executable = str(resolved)
    print(f"[INFO] 프로즌 모드 시뮬레이션: sys.executable={sys.executable}")

    engine = resolved.parent / "engine"
    soffice = office_mod.find_soffice()
    check("자동 탐색: soffice (env 미사용)", soffice is not None and str(engine) in soffice, soffice)
    java = hwp_mod._java()
    check("자동 탐색: java (env 미사용)", java is not None and str(engine) in java, java)
    cp = hwp_mod._classpath()
    check("자동 탐색: hwp classpath (env 미사용)", cp is not None and str(engine) in cp, cp)
    ffmpeg = video_mod.find_ffmpeg()
    check("자동 탐색: ffmpeg (env 미사용)", ffmpeg is not None and str(engine) in ffmpeg, ffmpeg)
    ffprobe = video_mod.find_ffprobe()
    check("자동 탐색: ffprobe (env 미사용)", ffprobe is not None and str(engine) in ffprobe, ffprobe)

    from app.version import current_version
    version = current_version()
    check("자동 탐색: 번들 VERSION 파일 (업데이트 확인 기능용)",
          version != "0.0.0", version)


def main():
    argv = sys.argv[1:]
    skip_hwp = "--skip-hwp" in argv
    frozen_exe = None
    if "--frozen-exe" in argv:
        frozen_exe = argv[argv.index("--frozen-exe") + 1]
    positional = [a for a in argv if not a.startswith("--") and a != frozen_exe]
    hwp_sample = Path(positional[0]) if positional else None

    if frozen_exe:
        simulate_frozen(frozen_exe)

    tmp = Path(tempfile.mkdtemp())
    try:
        # 1) DOCX -> PDF (번들 LibreOffice 실경로 검증 — DEC-002)
        #    글꼴을 지정하지 않은, 실사용자 문서와 동일한 조건의 DOCX를 그대로
        #    변환한다 — 여기서는 텍스트 정확도까지만 검증한다. 글꼴 미지정 문서의
        #    최종 렌더링은 호스트에 설치된 글꼴에 좌우되며 우리가 보장할 수 있는
        #    범위 밖이다(잔여 리스크, 아래 4번 참고).
        from docx import Document
        src_docx = tmp / "smoke.docx"
        doc = Document()
        doc.add_paragraph("LibreOffice 번들 스모크 테스트 — 드문 자모: 뷁 밟 닳 넋 앎 옳")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "항목"
        table.cell(0, 1).text = "값"
        table.cell(1, 0).text = "결과"
        table.cell(1, 1).text = "성공"
        doc.save(src_docx)

        try:
            out = converters.convert(src_docx, "pdf", tmp)
            check("DOCX→PDF 변환 완료", out.exists())
            verify_pdf_content(out, "뷁 밟 닳 넋 앎 옳")
        except ConversionError as e:
            check("DOCX→PDF 변환 완료", False, f"{e.key}: {e.detail}")
        except Exception as e:
            check("DOCX→PDF 변환 완료", False, f"예상 밖 예외 {type(e).__name__}: {e}")

        # 2) PPTX -> PDF (DEC-016 — office_to_pdf 경로 재사용)
        #    글꼴 미지정 — 실사용자 문서와 동일 조건, 텍스트 정확도까지만 검증
        #    (위 1번 DOCX와 같은 이유로 폰트 내장은 여기서 보장 대상이 아니다).
        try:
            from pptx import Presentation
            src_pptx = tmp / "smoke.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "PPTX 스모크 테스트"
            slide.placeholders[1].text = "드문 자모: 뷁 밟 닳 넋 앎 옳"
            prs.save(src_pptx)

            out = converters.convert(src_pptx, "pdf", tmp)
            check("PPTX→PDF 변환 완료", out.exists())
            verify_pdf_content(out, "뷁 밟 닳 넋 앎 옳")
        except ImportError:
            print("[SKIP] PPTX→PDF — python-pptx 없음(테스트 픽스처 전용 의존성)")
        except ConversionError as e:
            check("PPTX→PDF 변환 완료", False, f"{e.key}: {e.detail}")
        except Exception as e:
            check("PPTX→PDF 변환 완료", False, f"예상 밖 예외 {type(e).__name__}: {e}")

        # 2-1) PPTX에 Noto Sans KR을 명시 지정한 경우 — 이 경로는 우리가
        #      통제 가능하므로(사용자가 PPTX 안에서 이 폰트를 직접 골랐거나,
        #      향후 우리가 PPTX를 생성하게 될 경우와 동일 조건) 폰트 내장까지
        #      엄격 검증한다. python-pptx의 동아시아 글꼴(a:ea)도 DOCX의
        #      eastAsia와 동일한 OOXML 구조라 같은 방식으로 지정한다.
        try:
            from pptx import Presentation
            from pptx.oxml.ns import qn
            src_pptx2 = tmp / "smoke_font.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "폰트 지정 PPTX 테스트"
            tf = slide.placeholders[1].text_frame
            tf.text = "드문 자모: 뷁 밟 닳 넋 앎 옳"
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Noto Sans KR"
                    rpr = run._r.get_or_add_rPr()
                    ea = rpr.makeelement(qn("a:ea"), {"typeface": "Noto Sans KR"})
                    rpr.append(ea)
            prs.save(src_pptx2)

            out = converters.convert(src_pptx2, "pdf", tmp)
            check("PPTX(폰트 지정)→PDF 변환 완료", out.exists())
            verify_pdf_content(out, "뷁 밟 닳 넋 앎 옳", must_embed_font=True)
        except ImportError:
            print("[SKIP] PPTX(폰트 지정)→PDF — python-pptx 없음")
        except ConversionError as e:
            check("PPTX(폰트 지정)→PDF 변환 완료", False, f"{e.key}: {e.detail}")
        except Exception as e:
            check("PPTX(폰트 지정)→PDF 변환 완료", False, f"예상 밖 예외 {type(e).__name__}: {e}")

        # 3) 우리 자신이 생성하는 DOCX(PDF→DOCX·HWP→DOCX 출력, docx_build.py)는
        #    Noto Sans KR을 항상 명시적으로 지정한다(DEC-015) — 이 경로는 100%
        #    우리 통제 안에 있으므로, 그 결과를 다시 PDF로 렌더링해 번들 폰트가
        #    실제로 내장되는지까지 엄격하게 검증한다.
        from app.converters.docx_build import blocks_to_docx
        own_docx = blocks_to_docx(
            [{"type": "p", "text": "자체 생성 DOCX 검증 — 드문 자모: 뷁 밟 닳 넋 앎 옳"}],
            tmp / "own.docx",
        )
        try:
            out = converters.convert(own_docx, "pdf", tmp)
            check("자체 생성 DOCX→PDF 변환 완료", out.exists())
            verify_pdf_content(out, "뷁 밟 닳 넋 앎 옳", must_embed_font=True)
        except ConversionError as e:
            check("자체 생성 DOCX→PDF 변환 완료", False, f"{e.key}: {e.detail}")
        except Exception as e:
            check("자체 생성 DOCX→PDF 변환 완료", False, f"예상 밖 예외 {type(e).__name__}: {e}")

        # 4) HWP -> PDF (DEC-007 파이프라인: hwplib 구조 추출 → DOCX → LibreOffice)
        #    중간 산출물이 blocks_to_docx를 거치므로 여기도 폰트 내장을 엄격 검증한다.
        if not skip_hwp:
            check("HWP 샘플 존재", hwp_sample is not None and hwp_sample.exists(),
                  str(hwp_sample))
            try:
                out = converters.convert(hwp_sample, "pdf", tmp)
                check("HWP→PDF 변환 완료(전체 파이프라인)", out.exists())
                verify_pdf_content(out, "ABC", must_embed_font=True)  # 표.hwp 알려진 내용
            except ConversionError as e:
                check("HWP→PDF 변환 완료(전체 파이프라인)", False, f"{e.key}: {e.detail}")
            except Exception as e:
                check("HWP→PDF 변환 완료(전체 파이프라인)", False,
                      f"예상 밖 예외 {type(e).__name__}: {e}")

        # 5) DOCX -> HWP (DEC-017: hwplib 신규 생성 — 문단 텍스트 보존, 표는
        #    " | " 텍스트로 단순화). 생성된 HWP를 다시 TXT로 읽어 왕복 검증한다.
        if not skip_hwp:
            try:
                out = converters.convert(src_docx, "hwp", tmp)
                check("DOCX→HWP 변환 완료", out.exists())
                back = converters.convert(out, "txt", tmp)
                back_text = back.read_text(encoding="utf-8")
                check("DOCX→HWP: 문단 텍스트 보존", "뷁 밟 닳 넋 앎 옳" in back_text, repr(back_text[:120]))
                check("DOCX→HWP: 표→텍스트 단순화 내용 보존", "결과" in back_text and "성공" in back_text)
            except ConversionError as e:
                check("DOCX→HWP 변환 완료", False, f"{e.key}: {e.detail}")
            except Exception as e:
                check("DOCX→HWP 변환 완료", False, f"예상 밖 예외 {type(e).__name__}: {e}")

        # 6) 영상 -> MP4 (DEC-024: H.264 스트림은 재인코딩 없이 그대로 복사 —
        #    이 LGPL 빌드는 GPL 인코더가 없어 스스로 H.264를 만들 수 없으므로,
        #    저장소에 미리 넣어 둔 실제 샘플로 "카피" 경로만 검증한다).
        video_sample = REPO / "sidecar" / "ffmpeg" / "sample_h264_aac.mov"
        try:
            check("영상 샘플 존재", video_sample.exists(), str(video_sample))
            out = converters.convert(video_sample, "mp4", tmp)
            check("영상→MP4 변환 완료", out.exists())
            src_bytes = video_sample.read_bytes()
            check("영상→MP4: 결과 파일 비어있지 않음", out.stat().st_size > 0)
            check("영상→MP4: 원본보다 극단적으로 작지 않음(재인코딩 없이 카피됐다는 방증)",
                  out.stat().st_size > len(src_bytes) * 0.5,
                  f"src={len(src_bytes)} out={out.stat().st_size}")
        except ConversionError as e:
            check("영상→MP4 변환 완료", False, f"{e.key}: {e.detail}")
        except Exception as e:
            check("영상→MP4 변환 완료", False, f"예상 밖 예외 {type(e).__name__}: {e}")

        # 7) 이미지 포맷 변환 (DEC-025: Pillow — PyInstaller 패키징 후에도 실제로
        #    묶여서 동작하는지가 핵심. 새 네이티브 의존성이라 엔진들과 같은
        #    수준으로 패키징 경로 검증이 필요함).
        try:
            from PIL import Image
            src_png = tmp / "smoke.png"
            Image.new("RGB", (8, 8), (10, 20, 30)).save(src_png)
            out = converters.convert(src_png, "jpg", tmp)
            check("이미지 PNG→JPG 변환 완료", out.exists())
            with Image.open(out) as result:
                check("이미지 결과 포맷이 JPEG", result.format == "JPEG", result.format)
        except ConversionError as e:
            check("이미지 PNG→JPG 변환 완료", False, f"{e.key}: {e.detail}")
        except Exception as e:
            check("이미지 PNG→JPG 변환 완료", False, f"예상 밖 예외 {type(e).__name__}: {e}")

        print("스모크 전체 통과")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    main()
