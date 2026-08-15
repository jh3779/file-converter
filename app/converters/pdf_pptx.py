"""PDF → PPTX 전용 변환기 — pdf.py의 공유 추출 프리미티브(_paragraph_candidates·
_iter_lines·_container_to_runs·_iter_visuals·_visual_to_dict)를 재사용해
python-pptx 셰이프 API로 슬라이드를 만든다. pdf_docx.py(DOCX 전용 OOXML
작성 어댑터)와 대칭인 구조 — 구조 감사(2026-08)에서 세 관심사(공유 추출·
DOCX 작성·PPTX 작성)가 pdf.py 한 파일에 섞여 있던 것을 분리했다.
"""
from pathlib import Path

from .base import ConversionError
from .pdf import _container_to_runs, _iter_lines, _iter_visuals, _paragraph_candidates, _visual_to_dict


def pdf_to_pptx(src: Path, tmpdir: Path) -> Path:
    """PDF → PPTX, 한 페이지=한 이미지로 뭉개지 않고 줄 단위로 위치를 재구성
    (DEC-030). 각 텍스트 줄을 원본과 같은 위치·크기의 개별 텍스트 상자로
    복원하고, 굵게 판정은 pdf_to_docx와 같은 폰트 이름 휴리스틱(pdf.py의
    _container_to_runs)을 그대로 재사용한다. 이미지(LTImage)·사각형(LTRect)·
    직선(LTLine)·일반 곡선(LTCurve)도 원본과 같은 위치·크기로 재구성한다
    (DEC-036) — 표 테두리는 대개 LTLine(또는 LTRect)로 그려지므로 이제
    실제로 옮겨진다.

    **알려진 한계(정직하게 문서화, 스파이크로 직접 확인)**: (1) LTCurve(사각형·
    직선이 아닌 일반 곡선, 베지어 등)는 제어점을 그대로 직선으로 이은 다각형
    으로 근사한다 — 정확한 곡률은 재현되지 않는다(드문 경우, 표/도형 대부분은
    LTRect·LTLine). (2) 벡터·텍스트 사이의 원래 z-순서(위/아래 겹침)는 완전히
    보존되지 않는다 — 도형·이미지를 먼저 그리고 텍스트를 그 위에 올리는
    "배경 레이어" 방식으로 단순화(텍스트가 항상 보이는 게 더 흔한 실사용
    요구). (3) 디코딩 실패한 이미지(드문 필터·손상 데이터)는 조용히 건너뛴다
    (텍스트 보존 우선 원칙과 동일, `pdfminer.image.ImageWriter`가 처리 못하는
    포맷). (4) 재구성에 항상 Noto Sans KR을 지정하므로(DEC-015와 같은 이유)
    원본 폰트와 글자 폭이 달라 드물게 줄바꿈이 살짝 밀릴 수 있다. (5) 기울임은
    한글 글꼴 대부분에 별도 이탤릭 글리프가 없어(CJK 타이포그래피 관행)
    감지되지 않는 경우가 흔함(pdf_to_docx와 동일한 제약).
    """
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import MSO_ANCHOR

    EMU_PER_PT = 12700
    image_dir = tmpdir / "_pdf_pptx_images"
    layout = _extract_pdf_layout(src, image_dir)
    if not layout:
        raise ConversionError("err.corrupted", "페이지 없음")

    prs = Presentation()
    prs.slide_width = Emu(round(layout[0]["width"] * EMU_PER_PT))
    prs.slide_height = Emu(round(layout[0]["height"] * EMU_PER_PT))
    blank_layout = prs.slide_layouts[6]

    for page in layout:
        slide = prs.slides.add_slide(blank_layout)
        page_h = page["height"]
        for visual in page["visuals"]:
            _add_visual_to_slide(slide, visual, page_h, EMU_PER_PT)
        for line in page["lines"]:
            x0, y0, x1, y1 = line["bbox"]
            left = Emu(round(x0 * EMU_PER_PT))
            top = Emu(round((page_h - y1) * EMU_PER_PT))
            width = Emu(max(round((x1 - x0) * EMU_PER_PT), 1))
            height = Emu(max(round((y1 - y0) * EMU_PER_PT), 1))
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p = tf.paragraphs[0]
            for r in line["runs"]:
                run = p.add_run()
                run.text = r["text"]
                run.font.bold = r["bold"]
                run.font.italic = r["italic"]
                if r["size"]:
                    run.font.size = Pt(r["size"])
                run.font.name = "Noto Sans KR"

    out = tmpdir / (src.stem + ".pptx")
    prs.save(out)
    return out


def _extract_pdf_layout(src: Path, image_dir: Path) -> list[dict]:
    """PDF → 페이지별 [{"width","height","lines":[{"bbox","runs"}],"visuals":[...]}]
    — pdf_to_pptx 전용, 줄 단위 위치(bbox)까지 필요해 문단 단위로 합치는
    pdf.py의 _extract_pdf_blocks_by_page와는 분리했다(서식 판정 로직 자체는
    pdf.py의 _container_to_runs를 그대로 재사용). image_dir은 임베디드
    이미지를 디코딩해 저장할 폴더(호출자의 tmpdir 하위 — 변환 종료 후
    자동 정리됨).

    pdf_to_docx는 정렬(align)·밑줄까지 다뤄야 해서(DEC-037, 이미지·표
    테두리·밑줄 반영 개선) 이 함수 대신 pdf_docx.py의
    _extract_pdf_line_layout()을 쓴다."""
    from pdfminer.high_level import extract_pages
    from pdfminer.image import ImageWriter
    from pdfminer.pdfdocument import PDFPasswordIncorrect
    from pdfminer.psparser import PSException

    image_dir.mkdir(parents=True, exist_ok=True)
    writer = ImageWriter(str(image_dir))

    pages = []
    try:
        for page in extract_pages(str(src)):
            lines = []
            for container in _paragraph_candidates(page):
                for line in _iter_lines(container):
                    runs = _container_to_runs(line)
                    if runs:
                        lines.append({"bbox": line.bbox, "runs": runs})
            visuals = [
                v for item in _iter_visuals(page)
                if (v := _visual_to_dict(item, writer, image_dir)) is not None
            ]
            pages.append({"width": page.width, "height": page.height, "lines": lines, "visuals": visuals})
    except PDFPasswordIncorrect:
        raise ConversionError("err.password")
    except (PSException, ValueError, OSError) as e:
        raise ConversionError("err.corrupted", str(e))
    return pages


def _add_visual_to_slide(slide, visual: dict, page_h: float, emu_per_pt: int):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Emu, Pt

    kind = visual["kind"]
    x0, y0, x1, y1 = visual["bbox"]

    if kind == "image":
        left = Emu(round(x0 * emu_per_pt))
        top = Emu(round((page_h - y1) * emu_per_pt))
        width = Emu(max(round((x1 - x0) * emu_per_pt), 1))
        height = Emu(max(round((y1 - y0) * emu_per_pt), 1))
        try:
            slide.shapes.add_picture(visual["path"], left, top, width, height)
        except Exception:
            pass  # 저장은 됐지만 python-pptx가 못 여는 포맷 — 조용히 생략
        return

    linewidth = max(visual.get("linewidth") or 0.75, 0.25)

    if kind == "line":
        lx0, ly0 = visual["p0"]
        lx1, ly1 = visual["p1"]
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(round(lx0 * emu_per_pt)), Emu(round((page_h - ly0) * emu_per_pt)),
            Emu(round(lx1 * emu_per_pt)), Emu(round((page_h - ly1) * emu_per_pt)),
        )
        if visual["stroke"]:
            conn.line.color.rgb = RGBColor(*visual["stroke"])
        conn.line.width = Pt(linewidth)
        return

    if kind == "rect":
        left = Emu(round(x0 * emu_per_pt))
        top = Emu(round((page_h - y1) * emu_per_pt))
        width = Emu(max(round((x1 - x0) * emu_per_pt), 1))
        height = Emu(max(round((y1 - y0) * emu_per_pt), 1))
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    else:  # curve — 다각형 근사(freeform)
        pts = visual["pts"]
        if len(pts) < 2:
            return
        start_x, start_y = pts[0][0] * emu_per_pt, (page_h - pts[0][1]) * emu_per_pt
        fb = slide.shapes.build_freeform(start_x=start_x, start_y=start_y, scale=1.0)
        fb.add_line_segments(
            [(px * emu_per_pt, (page_h - py) * emu_per_pt) for px, py in pts[1:]], close=True)
        shape = fb.convert_to_shape()

    if visual.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*visual["fill"])
    else:
        shape.fill.background()

    if visual.get("stroke"):
        shape.line.color.rgb = RGBColor(*visual["stroke"])
        shape.line.width = Pt(linewidth)
    else:
        shape.line.fill.background()
