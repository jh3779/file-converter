"""PDF 변환 — PDF→TXT (REQ-F-004). PDF→DOCX는 v0.2."""
from pathlib import Path

from .base import ConversionError


def pdf_to_txt(src: Path, tmpdir: Path) -> Path:
    from pdfminer.high_level import extract_text
    from pdfminer.pdfdocument import PDFPasswordIncorrect
    from pdfminer.psparser import PSException
    try:
        text = extract_text(str(src))
    except PDFPasswordIncorrect:
        raise ConversionError("err.password")
    except (PSException, ValueError, OSError) as e:
        raise ConversionError("err.corrupted", str(e))
    out = tmpdir / (src.stem + ".txt")
    out.write_text(text, encoding="utf-8")
    return out


def pdf_to_docx(src: Path, tmpdir: Path) -> Path:
    """PDF → DOCX, 줄 단위로 원본과 같은 절대 위치에 배치해 재구성한다
    (DEC-037) — pdf_to_pptx(DEC-030)의 줄 단위 재구성과 같은 원리를, DOCX의
    레거시 "프레임" 기능(`w:framePr` — 문단을 페이지 절대좌표에 고정, Word·
    LibreOffice 둘 다 지원)으로 구현한다. 이전에는 문단을 순서대로 이어붙여
    흐르게 하는 방식으로 단순화했었다(DEC-010) — 이제 원본과 시각적으로
    훨씬 가까운 배치를 낸다. HWP→DOCX는 아직 이 방식으로 바뀌지 않았다
    (HWP 쪽 줄 단위 위치 정보를 새로 뽑아야 해서 범위 밖 — hwp.py의
    docx 변환은 여전히 blocks_to_docx 기반 흐르는 문서).

    pdf2docx(PyMuPDF)는 AGPL이라 사용 금지(DEC-007과 동일한 라이선스 원칙).

    **트레이드오프(사용자 확인 후 채택, 정직하게 문서화)**: 각 줄이 독립된
    프레임이 되므로, 결과 DOCX는 일반적인 "이어서 타이핑하면 자연스럽게
    다음 줄로 흐르는" 문서가 아니다 — 문장 중간에 텍스트를 추가해 그 줄
    프레임의 폭을 넘기면 다음 줄과 자연스럽게 안 이어진다(PDF→PPTX와
    동일한 트레이드오프 — 위치 정확도와 자유 편집 사이의 근본적 상충).

    이미지·벡터 도형(표 테두리 등)도 pdf_to_pptx(DEC-036)와 같은
    _iter_visuals/_visual_to_dict로 뽑아 원래 위치·크기로 재구성한다
    (이미지·표 테두리 반영 개선) — python-docx엔 pptx의 셰이프 API 같은
    고수준 도형 삽입 기능이 없어(자유 배치 가능한 도형은 python-docx가
    지원 안 함), 이미 검증된 _set_frame_pr(줄 텍스트가 쓰는 것과 동일)로
    빈 문단을 이미지/도형 자리에 절대 위치시키고, 이미지는 그 문단 안에
    그림을 넣고, 사각형·직선은 문단 테두리(w:pBdr, DOCX가 지원하는 몇 안
    되는 절대 배치 가능한 그리기 수단)로 표현한다.

    **알려진 한계(정직하게 문서화)**: (1) 사각형이 아닌 곡선(LTCurve)은
    pdf_to_pptx처럼 다각형으로 정확히 근사하지 않고 bounding box 사각형
    테두리로 더 단순하게 근사한다 — python-docx가 pptx의 프리폼(자유 곡선)
    도형 API에 대응하는 기능이 없어 문단 테두리(사각형만 표현 가능)로
    타협했다(드문 경우, 표 테두리 대부분은 LTRect·LTLine). (2) 벡터·텍스트
    사이의 원래 z-순서는 완전히 보존되지 않는다(pdf_to_pptx와 동일한 이유 —
    도형을 먼저, 텍스트를 나중에 배치해 텍스트가 항상 위에 보이도록 단순화).
    (3) 디코딩 실패한 이미지는 조용히 건너뛴다(pdf_to_pptx와 동일).
    (4) 재구성에 항상 Noto Sans KR을 지정하므로(DEC-015와 같은 이유) 원본
    폰트와 글자 폭이 달라 드물게 줄바꿈이 살짝 밀릴 수 있다. (5) 기울임은
    한글 글꼴 대부분에 별도 이탤릭 글리프가 없어(CJK 타이포그래피 관행)
    감지되지 않는 경우가 흔함(pdf_to_pptx와 동일한 제약).
    """
    from docx import Document
    from docx.enum.section import WD_SECTION
    from docx.shared import Emu, Pt

    from .docx_build import _align_map, _apply_run_style, _set_font

    EMU_PER_PT = 12700
    image_dir = tmpdir / "_pdf_docx_images"
    layout = _extract_pdf_line_layout(src, image_dir)
    if not layout:
        raise ConversionError("err.corrupted", "페이지 없음")

    doc = Document()

    def _size_section(section, page):
        section.page_width = Emu(round(page["width"] * EMU_PER_PT))
        section.page_height = Emu(round(page["height"] * EMU_PER_PT))
        for attr in ("left_margin", "right_margin", "top_margin", "bottom_margin"):
            setattr(section, attr, Emu(0))

    _size_section(doc.sections[0], layout[0])
    current_size = (layout[0]["width"], layout[0]["height"])

    for page_index, page in enumerate(layout):
        if page_index > 0:
            page_size = (page["width"], page["height"])
            if page_size == current_size:
                doc.add_page_break()
            else:
                # 페이지 크기가 바뀌는 경우(스캔 첨부문서의 가로/세로 혼합 등)
                # — 단순 페이지 나눔만으로는 section의 page_width/height가
                # 첫 페이지 크기로 고정된 채라, 아래 vAnchor="page" 기준
                # y좌표(page_h - y1) 계산이 실제 렌더링 페이지 높이와 어긋나
                # 텍스트가 밀려 보이는 버그가 있었다(PR 콘텐츠 리뷰로 발견).
                # 크기가 실제로 바뀔 때만 새 섹션을 열어 그 페이지 실제
                # 크기로 맞춘다 — 흔한 동일 크기 다중 페이지는 기존과 동일하게
                # 단순 페이지 나눔을 그대로 쓴다.
                section = doc.add_section(WD_SECTION.NEW_PAGE)
                _size_section(section, page)
                current_size = page_size
        page_h = page["height"]
        for visual in page["visuals"]:
            _add_visual_to_docx(doc, visual, page_h)
        for line in page["lines"]:
            x0, y0, x1, y1 = line["bbox"]
            p = doc.add_paragraph()
            pf = p.paragraph_format
            pf.space_before = Pt(0)
            pf.space_after = Pt(0)
            pf.line_spacing = 1.0
            for run_dict in line["runs"]:
                run = p.add_run(run_dict["text"])
                _set_font(run)
                _apply_run_style(run, run_dict)
            align = _align_map().get(line.get("align"))
            if align is not None:
                p.alignment = align
            _set_frame_pr(p, x_pt=x0, y_pt=page_h - y1,
                          w_pt=max(x1 - x0, 1), h_pt=max(y1 - y0, 1))

    out = tmpdir / (src.stem + ".docx")
    doc.save(out)
    return out


def _set_frame_pr(paragraph, x_pt: float, y_pt: float, w_pt: float, h_pt: float):
    """문단을 페이지 절대좌표에 고정한다(`w:framePr`) — DrawingML 플로팅
    도형보다 훨씬 단순한 레거시 기능이지만 Word·LibreOffice 둘 다 지원하며,
    이 용도(줄 단위 위치 재현)엔 충분하다. 스파이크로 실측: 지정한 x/y와
    LibreOffice로 렌더링한 PDF의 실제 텍스트 위치가 2pt 이내로 일치함을
    확인(hAnchor/vAnchor="page" 기준, 페이지 여백은 0으로 맞춰 좌표계를
    pdf_to_pptx의 EMU 변환과 동일하게 단순화)."""
    from docx.oxml.ns import qn

    TWIPS_PER_PT = 20
    pPr = paragraph._p.get_or_add_pPr()
    frame_pr = pPr.makeelement(qn("w:framePr"), {
        qn("w:w"): str(max(round(w_pt * TWIPS_PER_PT), 1)),
        qn("w:h"): str(max(round(h_pt * TWIPS_PER_PT), 1)),
        qn("w:hRule"): "exact",
        qn("w:hAnchor"): "page",
        qn("w:vAnchor"): "page",
        qn("w:x"): str(round(x_pt * TWIPS_PER_PT)),
        qn("w:y"): str(round(y_pt * TWIPS_PER_PT)),
        qn("w:wrap"): "none",
    })
    pPr.append(frame_pr)


def _add_visual_to_docx(doc, visual: dict, page_h: float):
    """이미지·벡터 도형(사각형/직선/곡선) 하나를 원래 위치·크기의 빈
    문단으로 배치한다(pdf_to_docx, 이미지·표 테두리 반영 개선) — pptx의
    셰이프 API에 대응하는 기능이 python-docx엔 없어, 이미 검증된
    _set_frame_pr로 절대 위치를 잡고 내용물만 종류별로 다르게 채운다:
    이미지는 그림 run, 사각형·직선·곡선은 문단 테두리(w:pBdr)+채움
    (w:shd). 디코딩 실패 이미지는 pdf_to_pptx와 동일하게 조용히 건너뛴다
    (빈 위치 문단만 남기지 않도록 실패 시 문단 자체를 제거)."""
    from docx.shared import Pt

    x0, y0, x1, y1 = visual["bbox"]
    w_pt, h_pt = max(x1 - x0, 1), max(y1 - y0, 1)
    top_pt = page_h - y1

    if visual["kind"] == "image":
        p = doc.add_paragraph()
        pf = p.paragraph_format
        pf.space_before = Pt(0)
        pf.space_after = Pt(0)
        try:
            p.add_run().add_picture(visual["path"], width=Pt(w_pt), height=Pt(h_pt))
        except Exception:
            p._p.getparent().remove(p._p)
            return
        _set_frame_pr(p, x_pt=x0, y_pt=top_pt, w_pt=w_pt, h_pt=h_pt)
        return

    p = doc.add_paragraph()
    pf = p.paragraph_format
    pf.space_before = Pt(0)
    pf.space_after = Pt(0)
    linewidth_pt = max(visual.get("linewidth") or 0.75, 0.25)
    if visual["kind"] == "line":
        lx0, ly0 = visual["p0"]
        lx1, ly1 = visual["p1"]
        # 변 하나만 쓴다 — 위/아래(또는 좌/우) 두 변을 다 쓰면 얇은 프레임
        # 안에서 겹쳐 보여야 할 두 선이 육안으로 갈라져(이중선처럼) 보이는
        # 문제를 로컬에서 실제 LibreOffice 렌더링으로 확인해 수정.
        sides = ["right"] if abs(ly1 - ly0) > abs(lx1 - lx0) else ["bottom"]
    else:  # rect·curve(bounding box 사각형으로 근사)
        sides = ["top", "bottom", "left", "right"]
    if visual.get("stroke"):
        _set_paragraph_borders(p, sides, visual["stroke"], linewidth_pt)
    if visual.get("fill"):
        _set_paragraph_shading(p, visual["fill"])
    _set_frame_pr(p, x_pt=x0, y_pt=top_pt, w_pt=w_pt, h_pt=h_pt)


def _set_paragraph_borders(paragraph, sides: list[str], color_rgb: tuple[int, int, int], width_pt: float):
    """문단 테두리(`w:pBdr`)를 지정한 변에만 건다 — python-docx엔 표 셀
    테두리 API는 있어도 일반 문단 테두리 API가 없어 raw XML로 만든다.
    표 테두리(사각형)는 4변 전부, 직선은 방향에 맞는 변 1~2개만 받는다."""
    from docx.oxml.ns import qn

    EIGHTHS_PER_PT = 8
    pPr = paragraph._p.get_or_add_pPr()
    pBdr = pPr.makeelement(qn("w:pBdr"), {})
    sz = str(max(round(width_pt * EIGHTHS_PER_PT), 2))
    color = "%02X%02X%02X" % color_rgb
    for side in sides:
        pBdr.append(pBdr.makeelement(qn(f"w:{side}"), {
            qn("w:val"): "single",
            qn("w:sz"): sz,
            qn("w:space"): "0",
            qn("w:color"): color,
        }))
    pPr.append(pBdr)


def _set_paragraph_shading(paragraph, color_rgb: tuple[int, int, int]):
    """문단 배경 채움(`w:shd`) — LTRect의 채움색(fill)을 표현한다."""
    from docx.oxml.ns import qn

    pPr = paragraph._p.get_or_add_pPr()
    pPr.append(pPr.makeelement(qn("w:shd"), {
        qn("w:val"): "clear",
        qn("w:color"): "auto",
        qn("w:fill"): "%02X%02X%02X" % color_rgb,
    }))


def pdf_to_pptx(src: Path, tmpdir: Path) -> Path:
    """PDF → PPTX, 한 페이지=한 이미지로 뭉개지 않고 줄 단위로 위치를 재구성
    (DEC-030). 각 텍스트 줄을 원본과 같은 위치·크기의 개별 텍스트 상자로
    복원하고, 굵게 판정은 pdf_to_docx와 같은 폰트 이름 휴리스틱(_container_to_runs)을
    그대로 재사용한다. 이미지(LTImage)·사각형(LTRect)·직선(LTLine)·일반
    곡선(LTCurve)도 원본과 같은 위치·크기로 재구성한다(DEC-036) — 표 테두리는
    대개 LTLine(또는 LTRect)로 그려지므로 이제 실제로 옮겨진다.

    **알려진 한계(정직하게 문서화, 스파이크로 직접 확인)**: (1) LTCurve(사각형·
    직선이 아닌 일반 곡선, 베지어 등)는 제어점을 그대로 직선으로 이은 다각형
    으로 근사한다 — 정확한 곡률은 재현되지 않는다(드문 경우, 표/도형 대부분은
    LTRect·LTLine). (2) 벡터·텍스트 사이의 원래 z-순서(위/아래 겹침)는 완전히
    보존되지 않는다 — 도형·이미지를 먼저 그리고 텍스트를 그 위에 올리는
    "배경 레이어" 방식으로 단순화(텍스트가 항상 보이는 게 더 흔한 실사용
    요구). (3) 디코딩 실패한 이미지(드문 필터·손상 데이터)는 조용히 건너뛴다
    (텍스트 보존 우선 원칙과 동일, `pdfminer.image.ImageWriter`가 처리 못하는
    포맷). (4) 재구성에 항상 Noto Sans KR을 지정하므로(DEC-015와 같은 이유)
    원본 폰트와 글자 폭이 달라 드물게 줄바꿈이 살짝 밀릴 수 있다. (5) 기울임은
    한글 글꼴 대부분에 별도 이탤릭 글리프가 없어(CJK 타이포그래피 관행)
    감지되지 않는 경우가 흔함(pdf_to_docx와 동일한 제약).
    """
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import MSO_ANCHOR

    EMU_PER_PT = 12700
    image_dir = tmpdir / "_pdf_pptx_images"
    layout = _extract_pdf_layout(src, image_dir)
    if not layout:
        raise ConversionError("err.corrupted", "페이지 없음")

    prs = Presentation()
    prs.slide_width = Emu(round(layout[0]["width"] * EMU_PER_PT))
    prs.slide_height = Emu(round(layout[0]["height"] * EMU_PER_PT))
    blank_layout = prs.slide_layouts[6]

    for page in layout:
        slide = prs.slides.add_slide(blank_layout)
        page_h = page["height"]
        for visual in page["visuals"]:
            _add_visual_to_slide(slide, visual, page_h, EMU_PER_PT)
        for line in page["lines"]:
            x0, y0, x1, y1 = line["bbox"]
            left = Emu(round(x0 * EMU_PER_PT))
            top = Emu(round((page_h - y1) * EMU_PER_PT))
            width = Emu(max(round((x1 - x0) * EMU_PER_PT), 1))
            height = Emu(max(round((y1 - y0) * EMU_PER_PT), 1))
            box = slide.shapes.add_textbox(left, top, width, height)
            tf = box.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.TOP
            p = tf.paragraphs[0]
            for r in line["runs"]:
                run = p.add_run()
                run.text = r["text"]
                run.font.bold = r["bold"]
                run.font.italic = r["italic"]
                if r["size"]:
                    run.font.size = Pt(r["size"])
                run.font.name = "Noto Sans KR"

    out = tmpdir / (src.stem + ".pptx")
    prs.save(out)
    return out


def _iter_lines(container):
    """컨테이너의 직계 LTTextLine들을 낸다. 줄로 안 묶인 컨테이너(LTFigure
    등, _paragraph_candidates 참고)는 통째로 줄 하나처럼 취급한다."""
    from pdfminer.layout import LTTextLine

    found_line = False
    for item in container:
        if isinstance(item, LTTextLine):
            found_line = True
            yield item
    if not found_line:
        yield container


def _extract_pdf_layout(src: Path, image_dir: Path) -> list[dict]:
    """PDF → 페이지별 [{"width","height","lines":[{"bbox","runs"}],"visuals":[...]}]
    — pdf_to_pptx 전용, 줄 단위 위치(bbox)까지 필요해 문단 단위로 합치는
    _extract_pdf_blocks와는 분리했다(서식 판정 로직 자체는 _container_to_runs를
    그대로 재사용). image_dir은 임베디드 이미지를 디코딩해 저장할 폴더
    (호출자의 tmpdir 하위 — 변환 종료 후 자동 정리됨).

    pdf_to_docx는 이미지·도형을 다루지 않으므로(DEC-037, 범위 밖으로 정직하게
    문서화) 이 함수 대신 이미지 추출 없는 _extract_pdf_line_layout()을 쓴다."""
    from pdfminer.high_level import extract_pages
    from pdfminer.image import ImageWriter
    from pdfminer.pdfdocument import PDFPasswordIncorrect
    from pdfminer.psparser import PSException

    image_dir.mkdir(parents=True, exist_ok=True)
    writer = ImageWriter(str(image_dir))

    pages = []
    try:
        for page in extract_pages(str(src)):
            lines = []
            for container in _paragraph_candidates(page):
                for line in _iter_lines(container):
                    runs = _container_to_runs(line)
                    if runs:
                        lines.append({"bbox": line.bbox, "runs": runs})
            visuals = [
                v for item in _iter_visuals(page)
                if (v := _visual_to_dict(item, writer, image_dir)) is not None
            ]
            pages.append({"width": page.width, "height": page.height, "lines": lines, "visuals": visuals})
    except PDFPasswordIncorrect:
        raise ConversionError("err.password")
    except (PSException, ValueError, OSError) as e:
        raise ConversionError("err.corrupted", str(e))
    return pages


_ALIGN_TOL = 3.0  # pt — 글자 위치 반올림오차 허용치


def _detect_alignment(container, page_width: float) -> str | None:
    """문단(컨테이너) 안 줄들의 가로 위치(bbox)로 정렬을 추정한다(DEC-040) —
    판정 로직 자체는 _classify_alignment(순수 함수, 단위 테스트 대상)에
    있고, 여기서는 pdfminer 컨테이너에서 줄별 bbox만 뽑아 넘긴다."""
    boxes = [line.bbox for line in _iter_lines(container)]
    return _classify_alignment(boxes, page_width)


def _classify_alignment(boxes: list[tuple[float, float, float, float]], page_width: float) -> str | None:
    """줄별 bbox((x0,y0,x1,y1), ...) 목록으로 정렬을 추정하는 순수 판정
    로직(DEC-040) — pdfminer 객체에 의존하지 않아 hand-crafted bbox로
    직접 단위 테스트할 수 있다.

    PDF는 정렬 자체를 담지 않는다 — 글자마다 절대 좌표만 있어, 렌더링된
    좌표에서 역산하는 수밖에 없다. 판단 근거가 약하면(짧은 문단·애매한
    여백) None을 돌려줘 "align" 필드 자체를 생략한다 — 잘못 추정해 원래
    의도와 다른 정렬로 억지로 맞추는 것보다, 아무것도 안 하는(문서 기본
    정렬 유지) 쪽이 안전하다는 원칙(DEC-027의 "서식 불명 시 안전한 기본값"과
    같은 태도).

    한 줄짜리 문단: 좌우 여백이 비슷하면 가운데(페이지 중심 기준이라 문서의
    실제 여백 폭을 몰라도 판단 가능). **오른쪽 정렬도, 왼쪽 정렬도 한
    줄만으로는 확정하지 않는다** — "오른쪽 여백이 작다"는 게 문서의 실제
    오른쪽 여백이 원래 좁아서인지 정말 오른쪽 정렬이라서인지 한 줄만
    봐서는 구분할 근거가 없다(페이지 가장자리에 딱 붙어야만 판단 가능한
    값이 되어 버려 실사용 문서의 정상적인 여백을 가진 오른쪽 정렬을 대부분
    놓친다 — 로컬 검증 중 발견해 범위를 좁힘). 처음엔 "오른쪽으로 확정 안
    하니 왼쪽으로 명시하자"로 구현했었는데, 그러면 한 줄짜리 오른쪽 정렬
    날짜·서명·짧은 제목 같은 실사용 문서의 정상 사례를 "left"로 확정해
    버리는 다른 방향의 오분류가 생겼다(자동 리뷰로 발견,
    `_classify_alignment([(400,700,540,712)], 612)`로 재현 확인 — 왼쪽
    여백 400·오른쪽 여백 72로 명백히 오른쪽에 가까운데 "left"가 나왔음) —
    "왼쪽인지 오른쪽인지 한 줄로는 모른다"가 정확한 상태이므로 None(판단
    보류)이 맞다. 여러 줄 문단의 오른쪽 정렬은 아래처럼 다른 방식(줄마다
    오른쪽 끝이 서로 일치하는지, 문서의 실제 여백 값과 무관)으로 판단해
    이 문제가 없다.
    여러 줄 문단: 마지막 줄을 뺀 "본문 줄"이 최소 2개 있어야 오른쪽·양쪽
    정렬을 판단한다 — 본문 줄이 1개뿐이면(전체 2줄 문단) 그 한 줄의 오른쪽
    끝과 "일치 비교"할 대상이 자기 자신뿐이라 항상 참이 되어, 평범한 2줄
    왼쪽 정렬 문단이 양쪽 정렬로 잘못 분류되는 회귀가 있었다(자동 리뷰로
    발견, `_classify_alignment([(72,700,300,712),(72,680,500,692)], 612)`로
    재현 확인 — 2줄 왼쪽 정렬인데 "justify"가 나왔음). 본문 줄이 부족하면
    오른쪽·양쪽 판정은 포기하고 왼쪽·가운데만 확인한다(부족한 근거로
    확정하지 않는다는 이 함수의 기본 원칙과 일치). 왼쪽 끝만 일치하면
    "left"를 명시적으로 돌려준다 — 왼쪽·오른쪽 끝이 둘 다 일치하면 양쪽,
    오른쪽만 일치하면 오른쪽, 어느 쪽도 아니면 각 줄의 중심이 일치하는지로
    가운데 정렬만 추가로 확인한다."""
    if not boxes:
        return None

    if len(boxes) == 1:
        x0, _, x1, _ = boxes[0]
        left_margin = x0
        right_margin = page_width - x1
        if left_margin <= _ALIGN_TOL:
            return None  # 왼쪽 여백이 거의 없음 — 판단 근거 부족
        if abs(left_margin - right_margin) <= _ALIGN_TOL * 2:
            return "center"
        return None  # 왼쪽인지 오른쪽인지 한 줄만으로는 판단 안 함(위 설명)

    lefts = [b[0] for b in boxes]
    rights_body = [b[2] for b in boxes[:-1]]  # 마지막 줄은 양쪽 정렬이어도 보통 짧다
    left_consistent = (max(lefts) - min(lefts)) <= _ALIGN_TOL
    # 본문 줄이 최소 2개는 있어야 "일치"가 의미 있다 — 1개뿐이면 항상 자기
    # 자신과 같아 무조건 참이 되므로(2줄 문단 오분류 버그, 위 설명) 오른쪽·
    # 양쪽 판정 자체를 포기한다.
    right_consistent = len(rights_body) >= 2 and (max(rights_body) - min(rights_body)) <= _ALIGN_TOL

    if left_consistent and right_consistent:
        # 알려진 한계(자동 리뷰 지적, 재현 확인·의도적으로 안 고침): 모든
        # 줄의 폭이 우연히 똑같은 가운데 정렬 문단은 좌우 끝이 둘 다
        # 일치해 버려 여기서 "justify"로 잘못 나올 수 있다. "좌우 여백이
        # 페이지 중심 기준으로 대칭인지"를 추가 판별 신호로 써서 가운데로
        # 우선시키는 방안도 검토했으나, 좌우 여백이 똑같은(예: 1인치 표준
        # 여백) 페이지의 흔한 진짜 양쪽 정렬 문단까지 가운데로 오분류하는
        # 더 나쁜 회귀가 생김을 재현으로 확인해 채택하지 않았다(희귀한
        # 엣지 케이스 하나보다 흔한 케이스를 깨뜨리지 않는 쪽을 택함).
        return "justify"
    if left_consistent:
        return "left"  # 왼쪽 정렬 — 명시적으로 반환(생략하면 HWP 기본값인 양쪽 정렬로 해석됨)
    if right_consistent:
        return "right"
    centers = [(b[0] + b[2]) / 2 for b in boxes]
    if max(centers) - min(centers) <= _ALIGN_TOL:
        return "center"
    return None


def _iter_visuals(container):
    """레이아웃 트리를 재귀 순회하며 이미지·벡터 도형(LTImage/LTRect/LTLine/
    LTCurve)을 찾는다. 텍스트 추출(_paragraph_candidates)과는 별개의 순회다
    — 이미지는 LTFigure 안에 중첩돼 있고 선/사각형은 대개 페이지 최상위에
    바로 있어(로컬 스파이크로 확인), 같은 트리를 목적에 맞게 두 번 걷는 편이
    하나로 합치는 것보다 단순하다."""
    from pdfminer.layout import LTContainer, LTCurve, LTImage

    for item in container:
        if isinstance(item, LTImage):
            yield item
        elif isinstance(item, LTCurve):  # LTRect·LTLine도 LTCurve의 하위 클래스
            yield item
        if isinstance(item, LTContainer):
            yield from _iter_visuals(item)


def _pdf_color_to_rgb(color) -> tuple[int, int, int] | None:
    """PDF 색상값(그레이스케일 float, RGB/CMYK 튜플)을 (r,g,b) 0~255 정수
    튜플로 변환한다. 이름 있는 색상(문자열)은 로컬 재현 범위 밖이라 변환
    실패 시 None을 돌려주고 호출자가 기본값(선 없음/채움 없음)으로 처리한다
    — 색상 미표현이 도형 자체의 유실로 이어지지는 않는다."""
    if color is None:
        return None
    try:
        if isinstance(color, (int, float)):
            v = round(color * 255)
            return (v, v, v)
        if isinstance(color, (list, tuple)):
            if len(color) == 3:
                return tuple(round(c * 255) for c in color)
            if len(color) == 4:
                c, m, y, k = color
                return (
                    round(255 * (1 - c) * (1 - k)),
                    round(255 * (1 - m) * (1 - k)),
                    round(255 * (1 - y) * (1 - k)),
                )
    except (TypeError, ValueError):
        pass
    return None


def _visual_to_dict(item, writer, image_dir: Path) -> dict | None:
    from pdfminer.layout import LTImage, LTLine, LTRect

    if isinstance(item, LTImage):
        try:
            name = writer.export_image(item)
        except Exception:
            return None  # 디코딩 실패(드문 필터·손상 데이터) — 조용히 건너뜀
        return {"kind": "image", "bbox": item.bbox, "path": str(image_dir / name)}

    stroke = _pdf_color_to_rgb(item.stroking_color) if item.stroke else None
    fill = _pdf_color_to_rgb(item.non_stroking_color) if item.fill else None
    if isinstance(item, LTRect):
        return {"kind": "rect", "bbox": item.bbox, "stroke": stroke, "fill": fill,
                "linewidth": item.linewidth}
    if isinstance(item, LTLine):
        return {"kind": "line", "bbox": item.bbox, "p0": item.pts[0], "p1": item.pts[1],
                "stroke": stroke, "linewidth": item.linewidth}
    # 사각형·직선이 아닌 일반 LTCurve — 제어점을 직선으로 이은 다각형으로 근사.
    return {"kind": "curve", "bbox": item.bbox, "pts": list(item.pts),
            "stroke": stroke, "fill": fill, "linewidth": item.linewidth}


def _add_visual_to_slide(slide, visual: dict, page_h: float, emu_per_pt: int):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.util import Emu, Pt

    kind = visual["kind"]
    x0, y0, x1, y1 = visual["bbox"]

    if kind == "image":
        left = Emu(round(x0 * emu_per_pt))
        top = Emu(round((page_h - y1) * emu_per_pt))
        width = Emu(max(round((x1 - x0) * emu_per_pt), 1))
        height = Emu(max(round((y1 - y0) * emu_per_pt), 1))
        try:
            slide.shapes.add_picture(visual["path"], left, top, width, height)
        except Exception:
            pass  # 저장은 됐지만 python-pptx가 못 여는 포맷 — 조용히 생략
        return

    linewidth = max(visual.get("linewidth") or 0.75, 0.25)

    if kind == "line":
        lx0, ly0 = visual["p0"]
        lx1, ly1 = visual["p1"]
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(round(lx0 * emu_per_pt)), Emu(round((page_h - ly0) * emu_per_pt)),
            Emu(round(lx1 * emu_per_pt)), Emu(round((page_h - ly1) * emu_per_pt)),
        )
        if visual["stroke"]:
            conn.line.color.rgb = RGBColor(*visual["stroke"])
        conn.line.width = Pt(linewidth)
        return

    if kind == "rect":
        left = Emu(round(x0 * emu_per_pt))
        top = Emu(round((page_h - y1) * emu_per_pt))
        width = Emu(max(round((x1 - x0) * emu_per_pt), 1))
        height = Emu(max(round((y1 - y0) * emu_per_pt), 1))
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    else:  # curve — 다각형 근사(freeform)
        pts = visual["pts"]
        if len(pts) < 2:
            return
        start_x, start_y = pts[0][0] * emu_per_pt, (page_h - pts[0][1]) * emu_per_pt
        fb = slide.shapes.build_freeform(start_x=start_x, start_y=start_y, scale=1.0)
        fb.add_line_segments(
            [(px * emu_per_pt, (page_h - py) * emu_per_pt) for px, py in pts[1:]], close=True)
        shape = fb.convert_to_shape()

    if visual.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*visual["fill"])
    else:
        shape.fill.background()

    if visual.get("stroke"):
        shape.line.color.rgb = RGBColor(*visual["stroke"])
        shape.line.width = Pt(linewidth)
    else:
        shape.line.fill.background()


def _extract_pdf_blocks(src: Path) -> list[dict]:
    """PDF → 문단 블록(서식 포함, DEC-027). pdfminer의 레이아웃 트리를 재귀
    순회해 "문단 하나"로 볼 수 있는 단위마다 글자 단위로 훑어 굵게/기울임/
    크기가 바뀌는 지점마다 run을 새로 만든다.

    굵게/기울임 판정은 pdfminer가 넘겨주는 폰트 리소스 이름(LTChar.fontname,
    예: "Caladea-Bold")에 "Bold"/"Italic"/"Oblique" 문자열이 포함되는지 보는
    휴리스틱이다 — PDF 자체에 "이 글자가 굵다"는 명시적 플래그가 없어 폰트
    이름에 의존할 수밖에 없다. **로컬 검증 결과, 이 휴리스틱은 굵기별로
    폰트 파일이 실제로 분리돼 있을 때만(한글 폰트 포함, 예: NotoSansKR-Bold·
    AppleSDGothicNeo-Bold 등 — 실사용 문서 대부분이 이런 폰트를 씀) 정확히
    동작함을 확인했다.** 다만 두 가지 알려진 한계가 있다: (1) 문서에 동아시아
    글꼴을 명시하지 않아 렌더러가 임의의 대체 글꼴 하나로 뭉뚱그려 그리면
    굵기 정보 자체가 사라져 감지 불가(로컬 재현 확인), (2) 기울임은 한글
    글꼴 대부분이 별도 이탤릭 글리프가 없어(CJK 타이포그래피 관행) 애초에
    렌더러가 반영하지 않는 경우가 흔함 — 이 경우 우리 휴리스틱의 실패가
    아니라 원본 자체에 감지할 서식이 없는 것. **밑줄은 감지하지 않는다** —
    PDF는 밑줄을 폰트 속성이 아니라 별도의 벡터 선(그림)으로 그리는 경우가
    많아 이 휴리스틱(폰트 이름 기반)으로는 원천적으로 판별할 수 없다(모든
    PDF run은 항상 underline=False — HWP→DOCX만 밑줄을 지원, DEC-027).
    텍스트 보존은 서식 감지 실패와 무관하게 항상 보장한다(서식 불명 시
    bold=False/italic=False로 안전하게 처리).

    정렬(DEC-040)은 문단 안 줄들의 가로 위치(bbox)를 페이지 폭과 비교하는
    별도 휴리스틱(_detect_alignment)으로 추정한다 — 판단 근거가 약하면 "align"
    필드 자체를 생략한다(문서 기본 정렬을 그대로 둠, 잘못된 추정보다 안전).
    """
    from pdfminer.high_level import extract_pages
    from pdfminer.pdfdocument import PDFPasswordIncorrect
    from pdfminer.psparser import PSException

    blocks = []
    try:
        for page in extract_pages(str(src)):
            for container in _paragraph_candidates(page):
                runs = _container_to_runs(container)
                if runs:
                    block = {"type": "p", "runs": runs}
                    align = _detect_alignment(container, page.width)
                    if align:
                        block["align"] = align
                    blocks.append(block)
    except PDFPasswordIncorrect:
        raise ConversionError("err.password")
    except (PSException, ValueError, OSError) as e:
        raise ConversionError("err.corrupted", str(e))
    return blocks


def _extract_pdf_blocks_by_page(src: Path) -> list[dict]:
    """PDF → 문단 블록(평문, 페이지 경계 보존 — DEC-039). pdf_to_hwp 전용.

    기존 pdf_to_hwp는 pdf_to_txt()(pdfminer extract_text())로 문서 전체를
    한 문자열로 뽑은 뒤 빈 줄 기준으로만 문단을 나눴는데, extract_text()가
    페이지가 바뀌는 지점을 별도로 표시하지 않아 PDF 여러 페이지의 텍스트가
    HWP 안에서 페이지 구분 없이 한 페이지처럼 이어 붙어 보였다(외부 QA
    피드백으로 재현 확인). extract_pages()로 페이지 단위로 직접 순회해 각
    페이지 첫 문단에 pageBreakBefore를 표시한다(JsonToHwp.java가 이를 실제
    쪽 나눔으로 반영). 서식(굵게 등)은 이 경로에서 필요 없어
    _container_to_runs로 뽑은 run들의 텍스트만 이어붙인다. 정렬(DEC-040)은
    _detect_alignment로 함께 추정해 JsonToHwp.java에 전달한다.
    """
    from pdfminer.high_level import extract_pages
    from pdfminer.pdfdocument import PDFPasswordIncorrect
    from pdfminer.psparser import PSException

    blocks = []
    try:
        for page_index, page in enumerate(extract_pages(str(src))):
            first_on_page = True
            for container in _paragraph_candidates(page):
                text = "".join(r["text"] for r in _container_to_runs(container))
                # JsonToHwp.java는 text.trim().isEmpty()인 문단을 버린다
                # (sidecar/hwp/JsonToHwp.java) — 여기서 공백뿐인 컨테이너를
                # "의미 있는 첫 문단"으로 인정하면, pageBreakBefore가 그
                # 컨테이너에 붙었다가 JsonToHwp 쪽에서 통째로 버려져 실제
                # 첫 문단에는 쪽 나눔이 반영되지 않는 버그가 있었다(자동
                # 리뷰로 발견) — 같은 기준(strip 후 빈 문자열)으로 걸러
                # first_on_page를 소비하지 않게 한다.
                if not text.strip():
                    continue
                block = {"type": "p", "text": text}
                if page_index > 0 and first_on_page:
                    block["pageBreakBefore"] = True
                align = _detect_alignment(container, page.width)
                if align:
                    block["align"] = align
                blocks.append(block)
                first_on_page = False
    except PDFPasswordIncorrect:
        raise ConversionError("err.password")
    except (PSException, ValueError, OSError) as e:
        raise ConversionError("err.corrupted", str(e))
    return blocks


def _extract_pdf_line_layout(src: Path, image_dir: Path) -> list[dict]:
    """PDF → 페이지별 [{"width","height","lines":[{"bbox","runs","align"}],"visuals":[...]}] —
    pdf_to_docx 전용(DEC-037·이미지/표 테두리 반영 개선). pdf_to_pptx용
    _extract_pdf_layout()과 이미지·도형(visuals) 추출은 같은 헬퍼
    (_iter_visuals/_visual_to_dict)를 그대로 재사용한다 — 이제 DOCX도
    이미지·표 테두리를 반영하므로 두 함수가 다시 갈라질 이유는 "align"
    (문단 단위 정렬, PPTX 쪽은 안 씀)뿐이라 유지한다. 정렬(DEC-040)은
    컨테이너(원본 문단) 단위로 한 번만 판정해(_detect_alignment) 그
    문단에 속한 모든 줄에 같은 값을 붙인다 — 줄마다 독립된 프레임으로
    배치되지만(pdf_to_docx), 정렬 판정 자체는 여러 줄을 함께 봐야 하는
    문단 단위 신호이기 때문이다."""
    from pdfminer.high_level import extract_pages
    from pdfminer.image import ImageWriter
    from pdfminer.pdfdocument import PDFPasswordIncorrect
    from pdfminer.psparser import PSException

    image_dir.mkdir(parents=True, exist_ok=True)
    writer = ImageWriter(str(image_dir))

    pages = []
    try:
        for page in extract_pages(str(src)):
            lines = []
            for container in _paragraph_candidates(page):
                align = _detect_alignment(container, page.width)
                for line in _iter_lines(container):
                    runs = _container_to_runs(line)
                    if runs:
                        entry = {"bbox": line.bbox, "runs": runs}
                        if align:
                            entry["align"] = align
                        lines.append(entry)
            visuals = [
                v for item in _iter_visuals(page)
                if (v := _visual_to_dict(item, writer, image_dir)) is not None
            ]
            pages.append({"width": page.width, "height": page.height, "lines": lines, "visuals": visuals})
    except PDFPasswordIncorrect:
        raise ConversionError("err.password")
    except (PSException, ValueError, OSError) as e:
        raise ConversionError("err.corrupted", str(e))
    return pages


def _paragraph_candidates(obj):
    """레이아웃 트리를 재귀 순회하며 "문단 하나"로 취급할 컨테이너를 찾는다.

    페이지 최상위 텍스트는 pdfminer가 이미 LTTextContainer(문단)로 잘
    묶어주지만, Form XObject(LTFigure — 예: 벡터 그래픽 안에 그려진 텍스트,
    일부 PDF 생성기가 텍스트 박스를 XObject로 감싸는 경우) 내부는 그렇게
    묶이지 않고 LTChar가 그대로 흩어져 있다(로컬 재현 확인: extract_text()는
    이 텍스트를 잡아내는데 LTTextContainer만 훑는 순회는 조용히 놓침 —
    코드 리뷰로 발견). LTFigure를 재귀적으로 더 훑어 중첩된 텍스트까지
    문단 후보로 모은다.
    """
    from pdfminer.layout import LTAnno, LTChar, LTContainer, LTTextContainer

    if isinstance(obj, LTTextContainer):
        yield obj
        return
    if isinstance(obj, LTContainer):
        if any(isinstance(c, (LTChar, LTAnno)) for c in obj):
            yield obj  # LTFigure처럼 글자를 직접 담은(줄로 안 묶인) 컨테이너
        for child in obj:
            if not isinstance(child, (LTChar, LTAnno)):
                yield from _paragraph_candidates(child)


def _iter_chars(container):
    """컨테이너 안의 LTChar/LTAnno를 순서대로 낸다. LTTextContainer(문단)
    안의 LTTextLine(줄) 한 겹, 또는 LTFigure처럼 줄로 안 묶이고 글자를
    바로 담은 컨테이너 둘 다 처리한다."""
    from pdfminer.layout import LTAnno, LTChar, LTTextLine

    for item in container:
        if isinstance(item, LTTextLine):
            yield from _iter_chars(item)
        elif isinstance(item, (LTChar, LTAnno)):
            yield item


def _container_to_runs(container) -> list[dict]:
    """컨테이너(문단 또는 줄)를 글자 단위로 훑어 굵게/기울임/크기가 바뀌는
    지점마다 run을 새로 만든다(DEC-027, pdf_to_docx·pdf_to_pptx 공용).

    굵게/기울임 판정은 pdfminer가 넘겨주는 폰트 리소스 이름(LTChar.fontname,
    예: "Caladea-Bold")에 "Bold"/"Italic"/"Oblique" 문자열이 포함되는지 보는
    휴리스틱이다 — PDF 자체에 "이 글자가 굵다"는 명시적 플래그가 없어 폰트
    이름에 의존할 수밖에 없다. **로컬 검증 결과, 이 휴리스틱은 굵기별로
    폰트 파일이 실제로 분리돼 있을 때만(한글 폰트 포함, 예: NotoSansKR-Bold·
    AppleSDGothicNeo-Bold 등 — 실사용 문서 대부분이 이런 폰트를 씀) 정확히
    동작함을 확인했다.** 다만 두 가지 알려진 한계가 있다: (1) 문서에 동아시아
    글꼴을 명시하지 않아 렌더러가 임의의 대체 글꼴 하나로 뭉뚱그려 그리면
    굵기 정보 자체가 사라져 감지 불가(로컬 재현 확인), (2) 기울임은 한글
    글꼴 대부분이 별도 이탤릭 글리프가 없어(CJK 타이포그래피 관행) 애초에
    렌더러가 반영하지 않는 경우가 흔함 — 이 경우 우리 휴리스틱의 실패가
    아니라 원본 자체에 감지할 서식이 없는 것. **밑줄은 감지하지 않는다** —
    PDF는 밑줄을 폰트 속성이 아니라 별도의 벡터 선(그림)으로 그리는 경우가
    많아 이 휴리스틱(폰트 이름 기반)으로는 원천적으로 판별할 수 없다(모든
    PDF run은 항상 underline=False — HWP→DOCX만 밑줄을 지원, DEC-027).
    텍스트 보존은 서식 감지 실패와 무관하게 항상 보장한다(서식 불명 시
    bold=False/italic=False로 안전하게 처리).
    """
    from pdfminer.layout import LTChar

    runs = []
    cur_text = []
    cur_style = (False, False, None)  # (bold, italic, size) — 셋 중 하나라도 바뀌면 run 분리

    def flush():
        if cur_text:
            bold, italic, size = cur_style
            runs.append({"text": "".join(cur_text), "bold": bool(bold),
                         "italic": bool(italic), "underline": False,
                         "size": size, "color": None})

    for ch in _iter_chars(container):
        if isinstance(ch, LTChar):
            fontname = ch.fontname.lower()
            style = ("bold" in fontname, "italic" in fontname or "oblique" in fontname,
                     round(ch.size, 1))
            text = ch.get_text()
        else:
            # LTAnno(가상 문자 — 줄바꿈·자간 보정 등, 폰트 정보 없음): 현재
            # run에 그대로 이어붙이고 서식 전환은 트리거하지 않는다. PDF의
            # 줄바꿈은 문단 내부 개행일 뿐이라(문단 경계는 컨테이너 단위로
            # 이미 나뉨) 공백으로 정규화한다.
            text = ch.get_text()
            text = " " if text == "\n" else text
            style = cur_style

        if style != cur_style and cur_text:
            flush()
            cur_text = []
        cur_style = style
        cur_text.append(text)
    flush()

    if runs:
        runs[0]["text"] = runs[0]["text"].lstrip()
        runs[-1]["text"] = runs[-1]["text"].rstrip()
        runs = [r for r in runs if r["text"]]
    return runs


_PDF_IMAGE_PILLOW_FORMAT = {"png": "PNG", "jpg": "JPEG"}


def pdf_to_images(src: Path, tmpdir: Path, ext: str = "png") -> Path:
    """PDF → 페이지별 이미지(PNG 또는 JPG), 원본 파일명 폴더 안에 저장
    (DEC-026, JPG 옵션은 DEC-043 — 외부 QA 피드백: PNG만 지원하던 것을 확장).

    pypdfium2(Apache-2.0/BSD-3-Clause, permissive) 사용 — LibreOffice의
    `--convert-to png`는 다중 페이지 PDF에서 첫 페이지 1장만 내보내는 것을
    직접 확인해 채택하지 않았다(엔진 조사 기록: DEC-026). PDF 페이지 렌더링
    결과(`bitmap.to_pil()`)는 항상 알파 채널 없는 RGB라 JPEG 저장 시 다른
    이미지 변환(image.py)처럼 투명 배경을 흰색으로 합성하는 처리가 필요
    없음을 로컬에서 직접 확인했다.

    반환값이 폴더 경로라는 점에서 다른 컨버터와 다르다 — output.py의
    finalize()가 결과가 폴더인지 파일인지 자동으로 판단해 처리한다.
    """
    import pypdfium2 as pdfium
    fmt = _PDF_IMAGE_PILLOW_FORMAT[ext]
    try:
        doc = pdfium.PdfDocument(src)
    except pdfium.PdfiumError as e:
        key = "err.password" if "password" in str(e).lower() else "err.corrupted"
        raise ConversionError(key, str(e))
    try:
        n_pages = len(doc)
        if n_pages == 0:
            raise ConversionError("err.corrupted", "페이지 없음")
        out_dir = tmpdir / src.stem
        out_dir.mkdir(parents=True, exist_ok=True)
        width = len(str(n_pages))
        for i, page in enumerate(doc, start=1):
            bitmap = page.render(scale=2.0)
            try:
                bitmap.to_pil().save(out_dir / f"page_{i:0{width}d}.{ext}", format=fmt)
            finally:
                bitmap.close()
                page.close()
        return out_dir
    finally:
        doc.close()
