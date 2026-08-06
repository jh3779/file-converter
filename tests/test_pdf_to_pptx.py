"""PDF→PPTX 테스트 — DEC-030.

핵심 검증: 페이지 하나를 이미지 한 장으로 뭉개는 대신, 줄 단위로 위치(bbox)를
재구성해 슬라이드에 개별 텍스트 상자로 넣는지(진짜 편집 가능한 텍스트인지),
굵게 판정(폰트 이름 휴리스틱, pdf_to_docx와 동일 로직 재사용)이 동작하는지.
표 테두리·이미지 같은 비텍스트 요소는 이 파이프라인 범위 밖(문서화된 한계,
note.pdf_to_pptx 고지로 안내) — 이 테스트는 텍스트 레이어만 검증한다.
"""
import shutil
import tempfile
import unittest
from pathlib import Path

from app import converters
from app.converters.base import ConversionError
from pptx import Presentation


def _mini_pdf(path: Path, pages: list[list[tuple]]):
    """손으로 만든 최소 PDF. pages: 페이지별 [(text, x, y, font_key, size)] 목록.
    font_key "F1"=Helvetica(일반), "F2"=Helvetica-Bold(굵게 감지용,
    pdf.py의 fontname.lower()에 "bold" 포함 여부 휴리스틱과 맞춤)."""
    page_obj_start = 3
    n = len(pages)
    font1_num = page_obj_start + 2 * n
    font2_num = font1_num + 1
    kids, page_bodies, content_bodies = [], [], []
    for i, lines in enumerate(pages):
        page_idx = page_obj_start + 2 * i
        content_idx = page_idx + 1
        kids.append(f"{page_idx} 0 R")
        ops = []
        for text, x, y, font_key, size in lines:
            ops.append(f"BT /{font_key} {size} Tf {x} {y} Td ({text}) Tj ET")
        content = "\n".join(ops).encode()
        stream = b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream"
        page_bodies.append(
            f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents {content_idx} 0 R"
            f" /Resources << /Font << /F1 {font1_num} 0 R /F2 {font2_num} 0 R >> >> >>".encode())
        content_bodies.append(stream)
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        ("<< /Type /Pages /Kids [" + " ".join(kids) + f"] /Count {n} >>").encode(),
    ]
    for pb, cb in zip(page_bodies, content_bodies):
        objs.append(pb)
        objs.append(cb)
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    objs.append(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold >>")

    buf = b"%PDF-1.4\n"
    offsets = []
    for i, o in enumerate(objs, 1):
        offsets.append(len(buf))
        buf += f"{i} 0 obj\n".encode() + o + b"\nendobj\n"
    xref = len(buf)
    buf += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        buf += f"{off:010d} 00000 n \n".encode()
    buf += (b"trailer\n<< /Size " + str(len(objs) + 1).encode() +
            b" /Root 1 0 R >>\nstartxref\n" + str(xref).encode() + b"\n%%EOF")
    path.write_bytes(buf)


class TestPdfToPptx(unittest.TestCase):
    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_slide_count_matches_page_count(self):
        src = self.tmp / "multi.pdf"
        _mini_pdf(src, [
            [("Page one", 50, 700, "F1", 12)],
            [("Page two", 50, 700, "F1", 12)],
            [("Page three", 50, 700, "F1", 12)],
        ])
        out = converters.convert(src, "pptx", self.tmp)
        self.assertTrue(out.exists())
        prs = Presentation(out)
        self.assertEqual(len(prs.slides.__iter__.__self__._sldIdLst), 3)

    def test_text_content_preserved_as_editable_run(self):
        src = self.tmp / "single.pdf"
        _mini_pdf(src, [[("Hello world", 50, 700, "F1", 12)]])
        out = converters.convert(src, "pptx", self.tmp)
        prs = Presentation(out)
        slide = list(prs.slides)[0]
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        self.assertIn("Hello world", texts)

    def test_line_position_roughly_matches_pdf_coordinates(self):
        """줄 bbox(PDF 포인트, 원점 좌하단) → EMU(원점 좌상단) 변환이 맞는지 —
        페이지 높이(792pt)에서 텍스트 y좌표(700pt)를 뺀 위치 근방에 텍스트
        상자가 있어야 한다."""
        src = self.tmp / "pos.pdf"
        _mini_pdf(src, [[("Positioned", 100, 700, "F1", 12)]])
        out = converters.convert(src, "pptx", self.tmp)
        prs = Presentation(out)
        slide = list(prs.slides)[0]
        box = next(s for s in slide.shapes if s.has_text_frame)
        EMU_PER_PT = 12700
        expected_left = 100 * EMU_PER_PT
        self.assertAlmostEqual(box.left, expected_left, delta=EMU_PER_PT * 2)
        # top은 (페이지높이 - 줄의 y1) 근방 — 정확한 y1은 폰트 상승폭에 좌우되므로
        # 대략 700pt 부근(오차 허용)인지만 확인한다.
        expected_top_approx = (792 - 700) * EMU_PER_PT
        self.assertAlmostEqual(box.top, expected_top_approx, delta=EMU_PER_PT * 15)

    def test_bold_font_detected_via_fontname_heuristic(self):
        """pdf_to_docx(DEC-027)와 같은 휴리스틱 재사용 — Helvetica-Bold로 그려진
        줄은 run.font.bold=True로 반영돼야 한다."""
        src = self.tmp / "bold.pdf"
        _mini_pdf(src, [[
            ("Normal", 50, 700, "F1", 12),
            ("Bold", 50, 680, "F2", 12),
        ]])
        out = converters.convert(src, "pptx", self.tmp)
        prs = Presentation(out)
        slide = list(prs.slides)[0]
        bold_flags = {}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    bold_flags[run.text] = run.font.bold
        self.assertEqual(bold_flags.get("Normal"), False)
        self.assertEqual(bold_flags.get("Bold"), True)

    def test_slide_size_matches_pdf_page_size(self):
        src = self.tmp / "size.pdf"
        _mini_pdf(src, [[("x", 50, 700, "F1", 12)]])
        out = converters.convert(src, "pptx", self.tmp)
        prs = Presentation(out)
        EMU_PER_PT = 12700
        self.assertEqual(prs.slide_width, round(612 * EMU_PER_PT))
        self.assertEqual(prs.slide_height, round(792 * EMU_PER_PT))

    def test_corrupted_pdf_rejected(self):
        src = self.tmp / "broken.pdf"
        src.write_bytes(b"not a real pdf")
        with self.assertRaises(ConversionError) as ctx:
            converters.convert(src, "pptx", self.tmp)
        self.assertEqual(ctx.exception.key, "err.corrupted")


if __name__ == "__main__":
    unittest.main()
