"""v0.2 파이프라인 테스트 — HWP→DOCX(사이드카 필요 시 스킵), PDF→TXT/DOCX."""
import shutil
import subprocess
import tempfile
import unittest
from pathlib import Path

from app import converters

REPO = Path(__file__).resolve().parents[1]
HWP_SAMPLE = REPO / "spike" / "hwplib" / "repo" / "sample_hwp" / "basic" / "표.hwp"
HWP_DISTRIBUTION = REPO / "spike" / "hwplib" / "repo" / "sample_hwp" / "distribution.hwp"


def _run_linesegdebug(hwp_path: Path):
    """LineSegDebug(테스트 전용 디버그 도구, sidecar/hwp/LineSegDebug.java)를
    직접 실행해 문단별 (lineAlignCount, LineSegItem 개수, [(textStart, vpos), ...])를
    돌려준다. 텍스트 왕복만으로는 DEC-018의 핵심 실패 모드(레이아웃 캐시가
    문단 길이·순서와 무관하게 고정값으로 나오는 것)를 검증할 수 없어서
    hwplib로 직접 열어 구조를 확인해야 한다."""
    from app.converters import hwp as hwp_mod
    java = hwp_mod._java()
    cp = hwp_mod._classpath()
    proc = subprocess.run([java, "-cp", cp, "LineSegDebug", str(hwp_path)],
                           capture_output=True, text=True, timeout=30)
    assert proc.returncode == 0, proc.stderr
    rows = []
    for line in proc.stdout.strip().splitlines():
        idx, line_align_count, seg_count, items = line.split("\t")
        pairs = []
        if items:
            for item in items.split(","):
                start, vpos = item.split(":")
                pairs.append((int(start), int(vpos)))
        rows.append((int(idx), int(line_align_count), int(seg_count), pairs))
    return rows


def _mini_pdf_centered(path: Path):
    """가운데 정렬된 한 줄짜리 최소 PDF(DEC-040) — 페이지 폭 612pt 기준
    좌우 여백이 정확히 같도록 Helvetica 18pt 폭(pdfminer 내장 AFM 테이블,
    tests/test_format_fidelity.py의 _helvetica_width와 같은 원리)을 직접
    계산해 배치한다. 단일 줄 정렬 감지 중 유일하게 신뢰도 높은 경우가
    가운데 정렬이라(페이지 중심 기준이라 문서의 실제 여백 폭을 몰라도
    판단 가능 — pdf.py._classify_alignment 참고) 이걸로 검증한다."""
    from pdfminer.pdffont import FONT_METRICS
    _, metrics = FONT_METRICS["Helvetica"]
    text = "Centered Line"
    size = 18
    width = sum(metrics.get(ch, metrics.get(" ", 500)) for ch in text) / 1000.0 * size
    x = (612 - width) / 2
    content = f"BT /F1 {size} Tf 1 0 0 1 {x:.2f} 700 Tm ({text}) Tj ET".encode()
    stream = b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream"
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R"
        b" /Resources << /Font << /F1 5 0 R >> >> >>",
        stream,
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    buf = b"%PDF-1.4\n"
    offsets = []
    for i, o in enumerate(objs, 1):
        offsets.append(len(buf))
        buf += f"{i} 0 obj\n".encode() + o + b"\nendobj\n"
    xref = len(buf)
    buf += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        buf += f"{off:010d} 00000 n \n".encode()
    buf += (b"trailer\n<< /Size " + str(len(objs) + 1).encode() +
            b" /Root 1 0 R >>\nstartxref\n" + str(xref).encode() + b"\n%%EOF")
    path.write_bytes(buf)


def _mini_pdf(path: Path):
    content = b"BT /F1 18 Tf 40 700 Td (Hello Converter) Tj ET"
    stream = b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream"
    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R"
        b" /Resources << /Font << /F1 5 0 R >> >> >>",
        stream,
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    buf = b"%PDF-1.4\n"
    offsets = []
    for i, o in enumerate(objs, 1):
        offsets.append(len(buf))
        buf += f"{i} 0 obj\n".encode() + o + b"\nendobj\n"
    xref = len(buf)
    buf += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        buf += f"{off:010d} 00000 n \n".encode()
    buf += (b"trailer\n<< /Size " + str(len(objs) + 1).encode() +
            b" /Root 1 0 R >>\nstartxref\n" + str(xref).encode() + b"\n%%EOF")
    path.write_bytes(buf)


class Base(unittest.TestCase):
    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)


class TestPdf(Base):
    def test_pdf_to_txt(self):
        pdf = self.tmp / "mini.pdf"
        _mini_pdf(pdf)
        out = converters.convert(pdf, "txt", self.tmp)
        self.assertIn("Hello Converter", out.read_text(encoding="utf-8"))

    def test_pdf_to_docx(self):
        from docx import Document
        pdf = self.tmp / "mini.pdf"
        _mini_pdf(pdf)
        out = converters.convert(pdf, "docx", self.tmp)
        texts = [p.text for p in Document(out).paragraphs]
        self.assertTrue(any("Hello Converter" in t for t in texts))


def _find_soffice():
    from app.converters import office
    return office.find_soffice()


@unittest.skipUnless(_find_soffice(), "LibreOffice(soffice) 없음 — 로컬/Windows CI에서만 실행")
class TestOffice(Base):
    """DEC-016: PPTX→PDF는 DOCX→PDF와 동일한 office_to_pdf 경로를 재사용한다."""

    def test_docx_to_pdf(self):
        from docx import Document
        src = self.tmp / "한글.docx"
        doc = Document()
        doc.add_paragraph("한글 DOCX→PDF 테스트")
        doc.save(src)
        out = converters.convert(src, "pdf", self.tmp)
        self.assertEqual(out.read_bytes()[:5], b"%PDF-")

    def test_pptx_to_pdf(self):
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx 없음 — 테스트 전용 의존성(pip install python-pptx)")
        src = self.tmp / "한글.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "한글 PPTX→PDF 테스트"
        prs.save(src)
        out = converters.convert(src, "pdf", self.tmp)
        self.assertEqual(out.read_bytes()[:5], b"%PDF-")


@unittest.skipUnless(
    HWP_SAMPLE.exists() and shutil.which("java"),
    "hwplib 샘플/JDK 없음 — spike 빌드 후 실행 (RESULT.md)")
class TestHwp(Base):
    def test_hwp_to_docx_preserves_table(self):
        from docx import Document
        out = converters.convert(HWP_SAMPLE, "docx", self.tmp)
        tables = Document(out).tables
        self.assertTrue(tables)
        cells = [c.text for row in tables[0].rows for c in row.cells]
        self.assertTrue(any("ABC" in c for c in cells))

    def test_hwp_to_txt(self):
        out = converters.convert(HWP_SAMPLE, "txt", self.tmp)
        self.assertIn("ABC", out.read_text(encoding="utf-8"))

    def test_pdf_to_hwp_roundtrip(self):
        """DEC-023: PDF→HWP는 텍스트 추출 기반(DEC-010과 같은 원칙 — PDF
        자체가 표 구조를 안 담고 있어 표 전용 처리는 필요 없음)."""
        src = self.tmp / "mini.pdf"
        _mini_pdf(src)
        out = converters.convert(src, "hwp", self.tmp)
        self.assertTrue(out.exists())
        back_dir = self.tmp / "back"
        back_dir.mkdir()
        back = converters.convert(out, "txt", back_dir)
        self.assertIn("Hello Converter", back.read_text(encoding="utf-8"))

    def test_pdf_to_hwp_alignment_roundtrip(self):
        """DEC-040: PDF의 가운데 정렬 문단이 bbox 휴리스틱으로 감지돼 HWP
        ParaShape에 실제로 반영되는지 — HWP→DOCX로 다시 왕복해 확인한다."""
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        src = self.tmp / "centered.pdf"
        _mini_pdf_centered(src)
        out = converters.convert(src, "hwp", self.tmp)
        back_dir = self.tmp / "back"
        back_dir.mkdir()
        back = converters.convert(out, "docx", back_dir)
        centered = next(p for p in Document(back).paragraphs if "Centered Line" in p.text)
        self.assertEqual(centered.alignment, WD_ALIGN_PARAGRAPH.CENTER)

    @unittest.skipUnless(HWP_DISTRIBUTION.exists(), "distribution.hwp 샘플 없음")
    def test_distribution_protected_hwp_still_readable(self):
        """OQ-006: '배포용(복사방지)' 문서는 편집·인쇄 제한이지 텍스트 암호화가
        아니므로 일반 경로로 정상 추출된다 — err.password로 오분류하지 않는다."""
        out = converters.convert(HWP_DISTRIBUTION, "txt", self.tmp)
        text = out.read_text(encoding="utf-8")
        self.assertTrue(text.strip())
        self.assertNotIn("�", text)

    def test_docx_to_hwp_roundtrip(self):
        """DEC-017/DEC-028: DOCX→HWP는 문단 텍스트를 보존하고, 표는 실제 HWP
        표 컨트롤로 새로 생성한다(이전엔 " | "로 이어붙인 텍스트로
        단순화했으나 DEC-017의 "hwplib은 표를 못 만든다"는 전제가 틀렸음이
        확인돼 정정됨). 결과 HWP를 다시 hwp_to_txt로 읽어 왕복 검증한다
        (희귀 한글 자모·한자 포함) — 표 구조 자체의 검증은
        test_format_fidelity.py 쪽에 별도로 있음."""
        from docx import Document
        src = self.tmp / "한글.docx"
        doc = Document()
        doc.add_paragraph("뷁 밟 닳 넋 앎 옳 훑 흙 삵 값 넓 얹 앉 닭 없")
        doc.add_paragraph("大韓民國 韓國語 漢字 契約書 委任狀")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "이름"
        table.cell(0, 1).text = "부서"
        table.cell(1, 0).text = "김철수"
        table.cell(1, 1).text = "영업1팀"
        doc.save(src)

        out = converters.convert(src, "hwp", self.tmp)
        self.assertTrue(out.exists())

        back_dir = self.tmp / "back"
        back_dir.mkdir()
        back = converters.convert(out, "txt", back_dir)
        text = back.read_text(encoding="utf-8")
        self.assertIn("뷁 밟 닳 넋 앎 옳 훑 흙 삵 값 넓 얹 앉 닭 없", text)
        self.assertIn("大韓民國 韓國語 漢字 契約書 委任狀", text)
        self.assertIn("이름", text)
        self.assertIn("김철수", text)
        self.assertIn("영업1팀", text)

    def test_docx_to_hwp_alignment_roundtrip(self):
        """DEC-040: DOCX 문단에 직접 지정된 정렬(가운데/오른쪽/왼쪽/양쪽)이
        HWP의 실제 정렬 ParaShape으로 반영되고, 다시 DOCX로 왕복해도 그대로
        남는지 확인한다. 명시적으로 정렬을 지정 안 한 문단은 HWP 문서 기본
        정렬(양쪽 정렬, hwplib 실측 확인)로 나와야 한다 — DOCX에서는 다시
        "정렬 안 지정"이 아니라 명시적 JUSTIFY로 나옴(HWP 쪽 기본값을 그대로
        옮기는 것도 이 기능의 의도된 동작)."""
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        src = self.tmp / "정렬.docx"
        doc = Document()
        doc.add_paragraph("기본 정렬")
        p_center = doc.add_paragraph("가운데 정렬")
        p_center.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p_right = doc.add_paragraph("오른쪽 정렬")
        p_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        p_left = doc.add_paragraph("왼쪽 정렬 명시")
        p_left.alignment = WD_ALIGN_PARAGRAPH.LEFT
        doc.save(src)

        out = converters.convert(src, "hwp", self.tmp)
        back_dir = self.tmp / "back"
        back_dir.mkdir()
        back = converters.convert(out, "docx", back_dir)
        by_text = {p.text: p.alignment for p in Document(back).paragraphs}

        self.assertEqual(by_text["기본 정렬"], WD_ALIGN_PARAGRAPH.JUSTIFY)
        self.assertEqual(by_text["가운데 정렬"], WD_ALIGN_PARAGRAPH.CENTER)
        self.assertEqual(by_text["오른쪽 정렬"], WD_ALIGN_PARAGRAPH.RIGHT)
        self.assertEqual(by_text["왼쪽 정렬 명시"], WD_ALIGN_PARAGRAPH.LEFT)

    def test_docx_to_hwp_preserves_long_wrapped_paragraph(self):
        """긴 문단(여러 줄로 감싸질 정도)이 HWP 레이아웃 캐시(LineSeg) 계산
        누락으로 실제 뷰어에서 겹쳐 보이던 문제 — hwplib 샘플 문서(distribution.hwp)
        구조를 근거로 sidecar/hwp/JsonToHwp.java가 줄바꿈을 계산하도록 수정.
        여기서는 파이썬 쪽에서 확인 가능한 텍스트 보존만 검증(레이아웃 구조
        자체는 로컬 스파이크로 별도 확인함)."""
        from docx import Document
        long_text = (
            "이것은 여러 줄로 감싸질 만큼 긴 문단입니다. " * 8
            + "뷁 밟 닳 넋 앎 옳 — 문단 끝 희귀 자모."
        )
        src = self.tmp / "긴문단.docx"
        doc = Document()
        doc.add_paragraph(long_text)
        doc.add_paragraph("짧은 두 번째 문단.")
        doc.save(src)

        out = converters.convert(src, "hwp", self.tmp)
        back_dir = self.tmp / "back"
        back_dir.mkdir()
        back = converters.convert(out, "txt", back_dir)
        text = back.read_text(encoding="utf-8")
        self.assertIn(long_text, text)
        self.assertIn("짧은 두 번째 문단.", text)

    def test_docx_to_hwp_lineseg_structure_matches_real_documents(self):
        """코드 리뷰 지적: 텍스트 왕복 테스트만으로는 DEC-018의 핵심 실패
        모드(LineSegItem이 문단 길이와 무관하게 항상 1개, 세로 위치가 모든
        문단에서 0으로 고정)를 전혀 잡지 못한다 — hwplib로 직접 열어
        레이아웃 캐시 구조(LineSegDebug, 테스트 전용 도구)까지 검증한다."""
        from docx import Document
        long_text = "구조 검증용 긴 문단입니다. " * 10
        src = self.tmp / "긴문단2.docx"
        doc = Document()
        doc.add_paragraph(long_text)
        doc.add_paragraph("짧은 두 번째 문단")
        doc.save(src)

        out = converters.convert(src, "hwp", self.tmp)
        rows = _run_linesegdebug(out)
        self.assertEqual(len(rows), 2)

        _, line_align_count0, seg_count0, items0 = rows[0]
        vpos0 = [v for _, v in items0]
        self.assertGreater(line_align_count0, 1)  # 여러 줄로 감싸져야 함(예전엔 항상 1)
        self.assertEqual(line_align_count0, seg_count0)  # 실제 문서 표본에서 확인된 불변식
        self.assertEqual(vpos0, sorted(vpos0))  # 세로 위치가 줄마다 증가
        self.assertEqual(len(set(vpos0)), len(vpos0))  # 모두 다른 값(0 고정 아님)

        _, line_align_count1, _, items1 = rows[1]
        self.assertEqual(line_align_count1, 1)
        # 다음 문단은 이전 문단의 마지막 줄보다 아래에서 시작해야 함 —
        # 예전 버그는 모든 문단이 vpos=0에서 시작해 서로 겹치는 것으로 기록됐음.
        self.assertGreater(items1[0][1], vpos0[-1])

    def test_docx_to_hwp_preserves_numbered_and_bullet_markers(self):
        """코드 리뷰 지적: DOCX 자동 번호·불릿은 문단 텍스트(w:t)가 아니라
        numbering.xml 서식으로 화면에만 그려지므로, item.text만 추출하면
        눈에 보이는 마커가 조용히 사라진다 — docx_extract가 numbering.xml을
        해석해 마커를 문단 앞에 붙이는지 전체 파이프라인으로 검증."""
        from docx import Document
        src = self.tmp / "목록.docx"
        doc = Document()
        doc.add_paragraph("첫 번째 항목", style="List Number")
        doc.add_paragraph("두 번째 항목", style="List Number")
        doc.add_paragraph("불릿 항목", style="List Bullet")
        doc.save(src)

        out = converters.convert(src, "hwp", self.tmp)
        back_dir = self.tmp / "back"
        back_dir.mkdir()
        back = converters.convert(out, "txt", back_dir)
        text = back.read_text(encoding="utf-8")
        self.assertIn("1. 첫 번째 항목", text)
        self.assertIn("2. 두 번째 항목", text)
        self.assertIn("불릿 항목", text)
        self.assertIn("•", text)


if __name__ == "__main__":
    unittest.main()
