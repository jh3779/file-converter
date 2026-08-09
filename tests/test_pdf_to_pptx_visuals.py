"""PDF→PPTX 이미지·벡터 도형(사각형/직선) 추출 테스트 — DEC-036.

test_pdf_to_pptx.py는 텍스트 레이어만 다루는 손으로 만든 최소 PDF를 쓴다
(당시엔 표 테두리·이미지가 범위 밖이었음). 이 파일은 같은 "손으로 최소
PDF 바이트를 직접 조립" 방식을 content-stream 연산자(re/m·l/Do)까지
확장해, LTRect(사각형 채움)·LTLine(직선)·LTImage(이미지 XObject)가 실제로
슬라이드에 위치·크기·색상까지 맞게 재구성되는지 검증한다.
"""
import shutil
import tempfile
import unittest
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE

from app import converters


def _build_pdf(path: Path, content_ops: str, extra_resources: str = "",
                extra_objects: list[bytes] | None = None, width=612, height=792):
    """단일 페이지 최소 PDF. extra_objects는 오브젝트 번호 5부터 순서대로
    배정되며(페이지=3, content=4), extra_resources 문자열 안에서 그 번호로
    직접 참조한다(예: 이미지 XObject를 5 0 obj로 넣었으면 "/XObject << /Im1
    5 0 R >>")."""
    extra_objects = extra_objects or []
    content = content_ops.encode()
    stream = b"<< /Length " + str(len(content)).encode() + b" >>\nstream\n" + content + b"\nendstream"

    objs = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 {width} {height}] /Contents 4 0 R"
         f" /Resources << {extra_resources} >> >>").encode(),
        stream,
        *extra_objects,
    ]

    buf = b"%PDF-1.4\n"
    offsets = []
    for i, o in enumerate(objs, 1):
        offsets.append(len(buf))
        buf += f"{i} 0 obj\n".encode() + o + b"\nendobj\n"
    xref = len(buf)
    buf += b"xref\n0 " + str(len(objs) + 1).encode() + b"\n0000000000 65535 f \n"
    for off in offsets:
        buf += f"{off:010d} 00000 n \n".encode()
    buf += (b"trailer\n<< /Size " + str(len(objs) + 1).encode() +
            b" /Root 1 0 R >>\nstartxref\n" + str(xref).encode() + b"\n%%EOF")
    path.write_bytes(buf)


class TestPdfToPptxVisuals(unittest.TestCase):
    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())

    def tearDown(self):
        shutil.rmtree(self.tmp, ignore_errors=True)

    def test_filled_rect_extracted_with_position_and_color(self):
        """0 0 1 rg(파랑 채움) 100 600 200 50 re f — 채워진 사각형 하나."""
        src = self.tmp / "rect.pdf"
        _build_pdf(src, "0 0 1 rg 100 600 200 50 re f")
        out = converters.convert(src, "pptx", self.tmp)

        from pptx import Presentation
        prs = Presentation(out)
        slide = list(prs.slides)[0]
        shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        self.assertEqual(len(shapes), 1)
        shape = shapes[0]

        EMU_PER_PT = 12700
        self.assertAlmostEqual(shape.left, 100 * EMU_PER_PT, delta=EMU_PER_PT)
        self.assertAlmostEqual(shape.width, 200 * EMU_PER_PT, delta=EMU_PER_PT)
        self.assertAlmostEqual(shape.height, 50 * EMU_PER_PT, delta=EMU_PER_PT)
        self.assertEqual(shape.fill.fore_color.rgb, (0, 0, 255))

    def test_stroked_line_extracted_with_position_and_color(self):
        """1 0 0 RG(빨강 선) 100 600 m 300 600 l S — 수평 직선 하나."""
        src = self.tmp / "line.pdf"
        _build_pdf(src, "1 0 0 RG 2 w 100 600 m 300 600 l S")
        out = converters.convert(src, "pptx", self.tmp)

        from pptx import Presentation
        prs = Presentation(out)
        slide = list(prs.slides)[0]
        lines = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
        self.assertEqual(len(lines), 1)
        line = lines[0]

        EMU_PER_PT = 12700
        self.assertAlmostEqual(line.width, 200 * EMU_PER_PT, delta=EMU_PER_PT)
        self.assertEqual(line.line.color.rgb, (255, 0, 0))

    def test_image_xobject_extracted_and_embedded_as_picture(self):
        """가공 없는 2x2 RGB 원시 이미지 XObject를 /Im1로 참조해 (50,600)
        위치에 100x100pt 크기로 배치 — 결과 슬라이드에 PICTURE 도형으로
        나와야 한다(픽셀 값 자체는 검증하지 않음, 위치·존재만 확인)."""
        src = self.tmp / "image.pdf"
        raw_rgb = bytes([
            255, 0, 0,    0, 255, 0,
            0, 0, 255,    255, 255, 0,
        ])  # 2x2, top row: red,green / bottom row: blue,yellow
        image_obj = (
            b"<< /Type /XObject /Subtype /Image /Width 2 /Height 2 "
            b"/ColorSpace /DeviceRGB /BitsPerComponent 8 /Length " +
            str(len(raw_rgb)).encode() + b" >>\nstream\n" + raw_rgb + b"\nendstream"
        )
        _build_pdf(
            src,
            "q 100 0 0 100 50 600 cm /Im1 Do Q",
            extra_resources="/XObject << /Im1 5 0 R >>",
            extra_objects=[image_obj],
        )
        out = converters.convert(src, "pptx", self.tmp)

        from pptx import Presentation
        prs = Presentation(out)
        slide = list(prs.slides)[0]
        pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        self.assertEqual(len(pictures), 1)
        pic = pictures[0]

        EMU_PER_PT = 12700
        self.assertAlmostEqual(pic.left, 50 * EMU_PER_PT, delta=EMU_PER_PT)
        self.assertAlmostEqual(pic.width, 100 * EMU_PER_PT, delta=EMU_PER_PT)
        self.assertAlmostEqual(pic.height, 100 * EMU_PER_PT, delta=EMU_PER_PT)


if __name__ == "__main__":
    unittest.main()
